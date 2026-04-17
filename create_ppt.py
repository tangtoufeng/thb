#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
二年级朗诵介绍PPT生成器
《七色花》——十二个月的故事节选
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_background(slide, color):
    """添加背景色"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()

def add_title_text(slide, text, top, font_size=44, color=None):
    """添加标题文字"""
    left = Inches(0.5)
    width = Inches(12.333)
    height = Inches(1.2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.name = "微软雅黑"
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    return textbox

def add_body_text(slide, text, top, font_size=24, color=None, left=0.8, width=11.7):
    """添加正文文字"""
    left = Inches(left)
    width = Inches(width)
    height = Inches(4.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.5
    p.font.size = Pt(font_size)
    p.font.name = "微软雅黑"
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(50, 50, 50)
    
    return textbox

def add_poem_text(slide, text, top, font_size=22, color=None):
    """添加诗歌/朗诵文字（居中）"""
    left = Inches(1)
    width = Inches(11.333)
    height = Inches(5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.8
    p.font.size = Pt(font_size)
    p.font.name = "楷体"
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(60, 60, 60)
    
    return textbox

def add_decorative_elements(slide, color):
    """添加装饰元素"""
    # 左上角装饰圆
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.2), Inches(0.2), Inches(0.6), Inches(0.6)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = color
    circle1.line.fill.background()
    
    # 右下角装饰圆
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(12.5), Inches(6.7), Inches(0.6), Inches(0.6)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = color
    circle2.line.fill.background()

# ==================== 第1页：封面 ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_background(slide1, RGBColor(135, 206, 235))  # 天蓝色背景

# 装饰元素
add_decorative_elements(slide1, RGBColor(255, 215, 0))

# 主标题
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "七色花"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(72)
p.font.bold = True
p.font.name = "微软雅黑"
p.font.color.rgb = RGBColor(255, 255, 255)

# 副标题
subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "十二个月的故事"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(40)
p.font.name = "楷体"
p.font.color.rgb = RGBColor(255, 223, 100)

# 朗诵介绍
info_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.8))
tf = info_box.text_frame
p = tf.paragraphs[0]
p.text = "二年级朗诵介绍"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(28)
p.font.name = "微软雅黑"
p.font.color.rgb = RGBColor(255, 255, 255)

# ==================== 第2页：故事简介 ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, RGBColor(255, 248, 220))  # 浅黄色背景

# 标题
add_title_text(slide2, "故事简介", Inches(0.5), font_size=48, color=RGBColor(139, 69, 19))

# 简介内容
intro_text = """《七色花》是俄罗斯著名童话故事

十二个月的故事讲述了：

正月、二月、三月三位老人
用魔法杖改变季节的神奇故事

他们让严寒退去
让春风吹拂
让万物复苏"""

add_body_text(slide2, intro_text, Inches(1.5), font_size=26, color=RGBColor(80, 50, 30))

# ==================== 第3页：正月的故事 ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, RGBColor(176, 224, 230))  # 淡蓝色背景

# 标题
add_title_text(slide3, "正月爷爷", Inches(0.4), font_size=46, color=RGBColor(25, 25, 112))

# 正月内容
january_text = """正月用自己的冰杖敲了一敲大地，
就念起咒文：

"严寒，不要在禁地的松林里
再把树木冻得发出破裂声吧。
不要再咬去松树和白桦树的树皮吧！
你们已经把鸦雀冻得够可怜啦，
你们把人类的住所也冻得够冰冷啦！"

老头儿话音刚落，树林里也变得静寂无声。
树木不再因为严寒而发出碎裂的响声，
白雪也开始像大块柔软的棉花，成堆地掉下来。"""

add_poem_text(slide3, january_text, Inches(1.3), font_size=22, color=RGBColor(40, 40, 100))

# ==================== 第4页：二月的故事 ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, RGBColor(230, 230, 250))  # 淡紫色背景

# 标题
add_title_text(slide4, "二月爷爷", Inches(0.4), font_size=46, color=RGBColor(75, 0, 130))

# 二月内容
february_text = """二月敲着手杖，摇着胡子，低声唱道：

"风啊，暴风啊，飓风啊，
你们尽力地吹吧！
旋风啊，暴风啊，
在夜晚时刮起来吧！
你们在云里高声地嚎叫，
你们吹过大地。
让雪像白蛇一样地在田野里奔驰吧！"

当他才把话讲完，
潮湿的狂风就在树枝中间喧响起来。
雪花飞舞着，白色的旋风刮过大地。"""

add_poem_text(slide4, february_text, Inches(1.3), font_size=22, color=RGBColor(60, 0, 100))

# ==================== 第5页：三月的故事 ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, RGBColor(144, 238, 144))  # 淡绿色背景

# 标题
add_title_text(slide5, "三月孩子", Inches(0.4), font_size=46, color=RGBColor(34, 139, 34))

# 三月内容
march_text = """三月用他孩子的声音大笑着，
并且响亮地歌唱着：

"雪已经不再是从前的样子了
——它们已经在田野里发黑啦。
湖上的冰已经裂开，
就好像有人把它们敲碎一样。
云飞得更快，天变得更高，
麻雀在屋脊上更加愉快地鸣叫。
细径和小路一天一天地露出黑色来，
杨柳树上的白花，像银子一样地闪着光。"""

add_poem_text(slide5, march_text, Inches(1.3), font_size=22, color=RGBColor(20, 80, 20))

# ==================== 第6页：春天来了 ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, RGBColor(255, 182, 193))  # 浅粉色背景

# 标题
add_title_text(slide6, "春天来了", Inches(0.4), font_size=46, color=RGBColor(199, 21, 133))

# 三月下半部分内容
spring_text = """"奔流吧！小河啊！
高涨吧！水潭啊！
蚂蚁们，在冬天的寒气之后，
一起爬出来吧！
大熊穿过了森林里倒下的树木。
鸟儿们都开始唱歌。
白雪花也盛开啦。"""

add_poem_text(slide6, spring_text, Inches(1.8), font_size=26, color=RGBColor(150, 20, 100))

# ==================== 第7页：朗诵提示 ====================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7, RGBColor(255, 250, 205))  # 浅金黄色背景

# 标题
add_title_text(slide7, "朗诵小提示", Inches(0.5), font_size=48, color=RGBColor(184, 134, 11))

# 朗诵提示内容
tips_text = """1. 正月的咒文——声音低沉、缓慢，像老爷爷说话

2. 二月的歌唱——声音有力、有气势，像风吹过的感觉

3. 三月的欢笑——声音明亮、欢快，像小孩子唱歌

4. 读到春天来临时——语气越来越轻快、越来越开心

5. 注意停顿和换气，让听众能感受到季节的变化"""

add_body_text(slide7, tips_text, Inches(1.5), font_size=24, color=RGBColor(100, 80, 30))

# ==================== 第8页：结束页 ====================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8, RGBColor(135, 206, 235))  # 天蓝色背景

# 装饰元素
add_decorative_elements(slide8, RGBColor(255, 215, 0))

# 感谢文字
thanks_box = slide8.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.5))
tf = thanks_box.text_frame
p = tf.paragraphs[0]
p.text = "谢谢大家！"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(60)
p.font.bold = True
p.font.name = "微软雅黑"
p.font.color.rgb = RGBColor(255, 255, 255)

# 副标题
end_box = slide8.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.333), Inches(1))
tf = end_box.text_frame
p = tf.paragraphs[0]
p.text = "请欣赏《七色花》朗诵"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(32)
p.font.name = "楷体"
p.font.color.rgb = RGBColor(255, 223, 100)

# 保存PPT
prs.save('/workspace/七色花_十二个月的故事_朗诵介绍.pptx')
print("PPT已成功生成！")
print("文件位置: /workspace/七色花_十二个月的故事_朗诵介绍.pptx")
print("\n共生成8页幻灯片：")
print("1. 封面 - 七色花·十二个月的故事")
print("2. 故事简介")
print("3. 正月爷爷")
print("4. 二月爷爷")
print("5. 三月孩子")
print("6. 春天来了")
print("7. 朗诵小提示")
print("8. 结束页 - 谢谢大家")
