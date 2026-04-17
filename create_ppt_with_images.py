#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
二年级朗诵介绍PPT生成器（含图片元素）
《七色花》——十二个月的故事节选
适合二年级小朋友的设计风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_background(slide, color):
    """添加背景色"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()

def add_title_text(slide, text, top, font_size=44, color=None, left=0.5, width=12.333):
    """添加标题文字"""
    left = Inches(left)
    width = Inches(width)
    height = Inches(1.2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.name = "微软雅黑"
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    return textbox

def add_body_text(slide, text, top, font_size=22, color=None, left=0.8, width=11.7, height=4.5):
    """添加正文文字"""
    left = Inches(left)
    width = Inches(width)
    height = Inches(height)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.5
    p.font.size = Pt(font_size)
    p.font.name = "微软雅黑"
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(50, 50, 50)
    
    return textbox

def add_poem_text(slide, text, top, font_size=20, color=None, left=0.5, width=12.333):
    """添加诗歌/朗诵文字（居中）"""
    left = Inches(left)
    width = Inches(width)
    height = Inches(4.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.8
    p.font.size = Pt(font_size)
    p.font.name = "楷体"
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = RGBColor(60, 60, 60)
    
    return textbox

def add_circle(slide, left, top, size, color):
    """添加圆形装饰"""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    return circle

def add_star(slide, left, top, size, color):
    """添加星形装饰"""
    star = slide.shapes.add_shape(
        MSO_SHAPE.STAR_5_POINT, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    star.fill.solid()
    star.fill.fore_color.rgb = color
    star.line.fill.background()
    return star

def add_cloud(slide, left, top, width, height, color):
    """添加云朵形状"""
    cloud = slide.shapes.add_shape(
        MSO_SHAPE.CLOUD, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    cloud.fill.solid()
    cloud.fill.fore_color.rgb = color
    cloud.line.fill.background()
    return cloud

def add_tree(slide, left, top, width, height, trunk_color, leaf_color):
    """添加树木（树干+树冠组合）"""
    # 树冠
    crown = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height * 0.7)
    )
    crown.fill.solid()
    crown.fill.fore_color.rgb = leaf_color
    crown.line.fill.background()
    
    # 树干
    trunk = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(left + width * 0.4), 
        Inches(top + height * 0.5), 
        Inches(width * 0.2), 
        Inches(height * 0.5)
    )
    trunk.fill.solid()
    trunk.fill.fore_color.rgb = trunk_color
    trunk.line.fill.background()
    
    return crown, trunk

def add_snowflake(slide, left, top, size, color):
    """添加雪花装饰（六角星）"""
    snow = slide.shapes.add_shape(
        MSO_SHAPE.STAR_6_POINT, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    snow.fill.solid()
    snow.fill.fore_color.rgb = color
    snow.line.fill.background()
    return snow

def add_sun(slide, left, top, size, color):
    """添加太阳装饰"""
    sun = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    sun.fill.solid()
    sun.fill.fore_color.rgb = color
    sun.line.fill.background()
    return sun

def add_bird(slide, left, top, width, height, color):
    """添加小鸟形状"""
    bird = slide.shapes.add_shape(
        MSO_SHAPE.WAVE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bird.fill.solid()
    bird.fill.fore_color.rgb = color
    bird.line.fill.background()
    return bird

def add_flower_decoration(slide, left, top, size, color):
    """添加花朵装饰（用五角星代替）"""
    flower = slide.shapes.add_shape(
        MSO_SHAPE.STAR_5_POINT, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    flower.fill.solid()
    flower.fill.fore_color.rgb = color
    flower.line.fill.background()
    return flower

def add_image_placeholder(slide, left, top, width, height, label):
    """添加图片占位框"""
    # 外框
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(240, 240, 240)
    frame.line.color.rgb = RGBColor(180, 180, 180)
    frame.line.width = Pt(2)
    
    # 占位文字
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top + height/2 - 0.2), Inches(width), Inches(0.4)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.name = "微软雅黑"
    p.font.color.rgb = RGBColor(150, 150, 150)
    
    return frame

# ==================== 第1页：封面 ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide1, RGBColor(135, 206, 235))  # 天蓝色背景

# 添加装饰花朵
add_flower_decoration(slide1, 0.3, 0.3, 0.8, RGBColor(255, 182, 193))
add_flower_decoration(slide1, 12.2, 0.5, 0.7, RGBColor(255, 192, 203))
add_flower_decoration(slide1, 0.5, 6.5, 0.6, RGBColor(221, 160, 221))
add_flower_decoration(slide1, 11.8, 6.2, 0.9, RGBColor(173, 216, 230))

# 主标题
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "七色花"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(72)
p.font.bold = True
p.font.name = "微软雅黑"
p.font.color.rgb = RGBColor(255, 255, 255)

# 副标题
subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "十二个月的故事"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(40)
p.font.name = "楷体"
p.font.color.rgb = RGBColor(255, 223, 100)

# 朗诵介绍
info_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.8))
tf = info_box.text_frame
p = tf.paragraphs[0]
p.text = "二年级朗诵介绍"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(28)
p.font.name = "微软雅黑"
p.font.color.rgb = RGBColor(255, 255, 255)

# 底部装饰星星
for i in range(5):
    add_star(slide1, 2 + i * 2.5, 6.5, 0.4, RGBColor(255, 215, 0))

# ==================== 第2页：故事简介 ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, RGBColor(255, 248, 220))  # 浅黄色背景

# 标题
add_title_text(slide2, "故事简介", Inches(0.4), font_size=48, color=RGBColor(139, 69, 19))

# 装饰元素
add_flower_decoration(slide2, 0.3, 0.3, 0.6, RGBColor(255, 182, 193))
add_flower_decoration(slide2, 12.5, 0.3, 0.6, RGBColor(173, 216, 230))

# 左侧：图片占位区
add_image_placeholder(slide2, 0.5, 1.6, 4, 3.5, "【可以放故事书封面图】")

# 右侧：简介内容
intro_text = """《七色花》是俄罗斯
著名童话故事

十二个月的故事讲述了：
正月、二月、三月
三位老人用魔法杖
改变季节的神奇故事

他们让严寒退去
让春风吹拂
让万物复苏"""

add_body_text(slide2, intro_text, Inches(1.8), font_size=24, color=RGBColor(80, 50, 30), 
                left=5, width=7.5, height=4)

# 底部装饰云朵
add_cloud(slide2, 1, 6.3, 2, 0.8, RGBColor(255, 255, 255))
add_cloud(slide2, 10, 6.5, 2.5, 0.9, RGBColor(255, 255, 255))

# ==================== 第3页：正月爷爷 ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, RGBColor(176, 224, 230))  # 淡蓝色背景

# 标题
add_title_text(slide3, "正月爷爷", Inches(0.3), font_size=46, color=RGBColor(25, 25, 112))

# 装饰雪花
add_snowflake(slide3, 0.3, 0.3, 0.6, RGBColor(255, 255, 255))
add_snowflake(slide3, 12.5, 0.4, 0.5, RGBColor(255, 255, 255))
add_snowflake(slide3, 0.5, 6.5, 0.5, RGBColor(255, 255, 255))
add_snowflake(slide3, 12.3, 6.3, 0.6, RGBColor(255, 255, 255))

# 左侧：树木插图（冬天的树）
add_tree(slide3, 0.5, 2.5, 1.2, 1.8, RGBColor(139, 90, 43), RGBColor(200, 200, 200))
add_tree(slide3, 1.8, 2.8, 1.0, 1.5, RGBColor(139, 90, 43), RGBColor(220, 220, 220))

# 右侧正月内容
january_text = """正月用自己的冰杖敲了一敲大地，
就念起咒文：

"严寒，不要在禁地的松林里
再把树木冻得发出破裂声吧。
不要再咬去松树和白桦树的树皮吧！
你们已经把鸦雀冻得够可怜啦，
你们把人类的住所也冻得够冰冷啦！"

老头儿话音刚落，树林里也变得静寂无声。
树木不再因为严寒而发出碎裂的响声，
白雪也开始像大块柔软的棉花，成堆地掉下来。"""

add_poem_text(slide3, january_text, Inches(1.0), font_size=19, color=RGBColor(40, 40, 100), 
              left=3.5, width=9)

# ==================== 第4页：二月爷爷 ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, RGBColor(230, 230, 250))  # 淡紫色背景

# 标题
add_title_text(slide4, "二月爷爷", Inches(0.3), font_size=46, color=RGBColor(75, 0, 130))

# 装饰元素 - 风的线条
for i in range(3):
    add_circle(slide4, 0.3 + i * 0.4, 0.5 + i * 0.3, 0.3, RGBColor(200, 200, 255))
    add_circle(slide4, 12.3 + i * 0.3, 0.4 + i * 0.4, 0.25, RGBColor(200, 200, 255))

# 左侧：风雪插图占位
add_image_placeholder(slide4, 0.3, 1.5, 3.5, 3, "【可以放风雪图片】")

# 右侧二月内容
february_text = """二月敲着手杖，摇着胡子，低声唱道：

"风啊，暴风啊，飓风啊，
你们尽力地吹吧！
旋风啊，暴风啊，
在夜晚时刮起来吧！
你们在云里高声地嚎叫，
你们吹过大地。
让雪像白蛇一样地在田野里奔驰吧！"

当他才把话讲完，
潮湿的狂风就在树枝中间喧响起来。
雪花飞舞着，白色的旋风刮过大地。"""

add_poem_text(slide4, february_text, Inches(1.0), font_size=19, color=RGBColor(60, 0, 100),
              left=4, width=9)

# 底部雪花装饰
for i in range(4):
    add_snowflake(slide4, 1 + i * 3, 6.5, 0.4, RGBColor(255, 255, 255))

# ==================== 第5页：三月孩子 ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, RGBColor(144, 238, 144))  # 淡绿色背景

# 标题
add_title_text(slide5, "三月孩子", Inches(0.3), font_size=46, color=RGBColor(34, 139, 34))

# 装饰元素 - 太阳
add_sun(slide5, 0.3, 0.3, 0.8, RGBColor(255, 215, 0))

# 装饰小鸟
add_bird(slide5, 11.5, 0.5, 1.2, 0.5, RGBColor(100, 149, 237))
add_bird(slide5, 11.8, 0.8, 1.0, 0.4, RGBColor(70, 130, 180))

# 左侧：春天景色占位
add_image_placeholder(slide5, 0.3, 1.5, 3.5, 3, "【可以放春天图片】")

# 右侧三月内容
march_text = """三月用他孩子的声音大笑着，
并且响亮地歌唱着：

"雪已经不再是从前的样子了
——它们已经在田野里发黑啦。
湖上的冰已经裂开，
就好像有人把它们敲碎一样。
云飞得更快，天变得更高，
麻雀在屋脊上更加愉快地鸣叫。
细径和小路一天一天地露出黑色来，
杨柳树上的白花，像银子一样地闪着光。"""

add_poem_text(slide5, march_text, Inches(1.0), font_size=19, color=RGBColor(20, 80, 20),
              left=4, width=9)

# 底部花朵装饰
for i in range(5):
    add_flower_decoration(slide5, 1 + i * 2.5, 6.4, 0.5, RGBColor(255, 192, 203))

# ==================== 第6页：春天来了 ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, RGBColor(255, 182, 193))  # 浅粉色背景

# 标题
add_title_text(slide6, "春天来了", Inches(0.3), font_size=46, color=RGBColor(199, 21, 133))

# 装饰元素 - 大太阳
add_sun(slide6, 0.2, 0.2, 1.0, RGBColor(255, 215, 0))

# 装饰小鸟群
for i in range(3):
    add_bird(slide6, 11 + i * 0.5, 0.4 + i * 0.3, 0.8, 0.35, RGBColor(100, 149, 237))

# 左侧：春天动物占位
add_image_placeholder(slide6, 0.3, 1.5, 3.5, 3, "【可以放小动物图片】")

# 右侧三月下半部分内容
spring_text = """"奔流吧！小河啊！
高涨吧！水潭啊！
蚂蚁们，在冬天的寒气之后，
一起爬出来吧！
大熊穿过了森林里倒下的树木。
鸟儿们都开始唱歌。
白雪花也盛开啦。"""

add_poem_text(slide6, spring_text, Inches(1.8), font_size=24, color=RGBColor(150, 20, 100),
              left=4, width=9)

# 底部装饰
add_flower_decoration(slide6, 1, 6.5, 0.6, RGBColor(255, 255, 255))
add_flower_decoration(slide6, 3, 6.3, 0.5, RGBColor(255, 192, 203))
add_flower_decoration(slide6, 10, 6.4, 0.6, RGBColor(221, 160, 221))
add_flower_decoration(slide6, 11.8, 6.6, 0.5, RGBColor(173, 216, 230))

# ==================== 第7页：朗诵提示 ====================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7, RGBColor(255, 250, 205))  # 浅金黄色背景

# 标题
add_title_text(slide7, "朗诵小提示", Inches(0.4), font_size=48, color=RGBColor(184, 134, 11))

# 装饰元素
add_star(slide7, 0.3, 0.3, 0.5, RGBColor(255, 215, 0))
add_star(slide7, 12.5, 0.3, 0.5, RGBColor(255, 215, 0))

# 左侧：小朋友朗诵图片占位
add_image_placeholder(slide7, 0.5, 1.5, 4, 4, "【可以放小朋友朗诵照片】")

# 右侧朗诵提示内容
tips_text = """1. 正月的咒文
   ——声音低沉、缓慢
   像老爷爷说话

2. 二月的歌唱
   ——声音有力、有气势
   像风吹过的感觉

3. 三月的欢笑
   ——声音明亮、欢快
   像小孩子唱歌

4. 读到春天来临时
   ——语气越来越轻快
   越来越开心

5. 注意停顿和换气
   让听众能感受到
   季节的变化"""

add_body_text(slide7, tips_text, Inches(1.5), font_size=22, color=RGBColor(100, 80, 30),
              left=5, width=7.5, height=5)

# 底部装饰花朵
for i in range(6):
    add_flower_decoration(slide7, 1 + i * 2, 6.6, 0.4, RGBColor(255, 182, 193))

# ==================== 第8页：结束页 ====================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8, RGBColor(135, 206, 235))  # 天蓝色背景

# 装饰元素
add_flower_decoration(slide8, 0.3, 0.3, 0.8, RGBColor(255, 182, 193))
add_flower_decoration(slide8, 12.2, 0.5, 0.7, RGBColor(255, 192, 203))
add_flower_decoration(slide8, 0.5, 6.5, 0.6, RGBColor(221, 160, 221))
add_flower_decoration(slide8, 11.8, 6.2, 0.9, RGBColor(173, 216, 230))

# 感谢文字
thanks_box = slide8.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
tf = thanks_box.text_frame
p = tf.paragraphs[0]
p.text = "谢谢大家！"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(60)
p.font.bold = True
p.font.name = "微软雅黑"
p.font.color.rgb = RGBColor(255, 255, 255)

# 副标题
end_box = slide8.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
tf = end_box.text_frame
p = tf.paragraphs[0]
p.text = "请欣赏《七色花》朗诵"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(32)
p.font.name = "楷体"
p.font.color.rgb = RGBColor(255, 223, 100)

# 底部装饰星星
for i in range(5):
    add_star(slide8, 2 + i * 2.5, 6.2, 0.5, RGBColor(255, 215, 0))

# 保存PPT
prs.save('/workspace/七色花_十二个月的故事_儿童版.pptx')
print("🎉 PPT已成功生成！")
print("=" * 50)
print("📁 文件位置: /workspace/七色花_十二个月的故事_儿童版.pptx")
print("=" * 50)
print("\n📊 共生成8页幻灯片：")
print("   1. 封面 - 七色花·十二个月的故事")
print("   2. 故事简介 - 带图片占位区")
print("   3. 正月爷爷 - 带树木、雪花装饰")
print("   4. 二月爷爷 - 带风雪图片占位区")
print("   5. 三月孩子 - 带太阳、小鸟、花朵装饰")
print("   6. 春天来了 - 带小动物图片占位区")
print("   7. 朗诵小提示 - 带小朋友照片占位区")
print("   8. 结束页 - 谢谢大家")
print("=" * 50)
print("\n✨ 特色功能：")
print("   • 添加了星星、雪花、太阳、小鸟等可爱装饰")
print("   • 设置了多个图片占位区，方便插入真实图片")
print("   • 配色丰富，适合二年级小朋友")
print("   • 字体大小适中，易于阅读")
